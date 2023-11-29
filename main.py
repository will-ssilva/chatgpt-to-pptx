import os
import re
import openai
import tkinter as tk
from pptx.util import Inches
from pptx import Presentation
from image_crawler import ImageCrawler
from message import message

global root
root = Presentation('modelo.pptx')

def delete_all_slides():
    for i in range(len(root.slides) - 1, -1, -1):
        r_id = root.slides._sldIdLst[i].rId
        root.part.drop_rel(r_id)
        del root.slides._sldIdLst[i]

def create_title_slide(title, subtitle):
    layout = root.slide_layouts[0]
    slide = root.slides.add_slide(layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def create_section_header_slide(title):
    layout = root.slide_layouts[2]
    slide = root.slides.add_slide(layout)
    slide.shapes.title.text = title

def create_title_and_content_slide(title, content):
    layout = root.slide_layouts[1]
    slide = root.slides.add_slide(layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

def create_title_and_content_and_image_slide(title, content, image_query):
    
    layout = root.slide_layouts[8]
    slide = root.slides.add_slide(layout)
    slide.shapes.title.text = title
    slide.placeholders[2].text = content
    img_path = ImageCrawler().google_crawler(image_query)
    picture = slide.placeholders[1].insert_picture(img_path)
    """picture.left = Inches(3)
    picture.top = Inches(0.5)
    picture.width = Inches(4)
    picture.height = Inches(3)"""

def find_text_in_between_tags(text, start_tag, end_tag):
    start_pos = text.find(start_tag)
    end_pos = text.find(end_tag)
    result = []
    while start_pos > -1 and end_pos > -1:
        text_between_tags = text[start_pos + len(start_tag):end_pos]
        result.append(text_between_tags)
        start_pos = text.find(start_tag, end_pos + len(end_tag))
        end_pos = text.find(end_tag, start_pos)
    res1 = "".join(result)
    res2 = re.sub(r"\[IMAGE\].*?\[/IMAGE\]\n", '', res1)
    if len(result) > 0:
        return res2
    else:
        return ""

def search_for_slide_type(text):
    tags = ["[L_TS]", "[L_CS]", "[L_IS]", "[L_THS]"]
    found_text = next((string for string in tags if string in text), None)
    return found_text

def create_slide_type(slide_type, slide):
    slide = str(slide)

    if slide_type == "[L_TS]":
        string_title = find_text_in_between_tags(slide, "[TITLE]", "[/TITLE]")
        string_subtitle = find_text_in_between_tags(slide, "[SUBTITLE]", "[/SUBTITLE]")
        create_title_slide(string_title, string_subtitle)

    elif slide_type == "[L_CS]":
        string_title = find_text_in_between_tags(slide, "[TITLE]", "[/TITLE]")
        string_content = find_text_in_between_tags(slide, "[CONTENT]", "[/CONTENT]")
        create_title_and_content_slide("".join(string_title), "".join(string_content))

    elif slide_type == "[L_IS]":
        string_title = find_text_in_between_tags(slide, "[TITLE]", "[/TITLE]")
        string_content = find_text_in_between_tags(slide, "[CONTENT]", "[/CONTENT]")
        string_image = find_text_in_between_tags(slide, "[IMAGE]", "[/IMAGE]")
        create_title_and_content_and_image_slide("".join(string_title), "".join(string_content), "".join(string_image))

    elif slide_type == "[L_THS]":
        string_title = find_text_in_between_tags(slide, "[TITLE]", "[/TITLE]")
        create_section_header_slide("".join(string_title))

def find_title():
    return root.slides[0].shapes.title.text


def parse_response(reply):
    delete_all_slides()
    list_of_slides = reply.split("[SLIDEBREAK]")
    for slide_string in list_of_slides:
        slide_type = search_for_slide_type(slide_string)
        create_slide_type(slide_type, slide_string)

    slide_file_title = find_title()
    root.save(f"{slide_file_title}.pptx")

    return rf"Feito! {slide_file_title} está pronto! Você pode encontrá-lo em {os.getcwd()}\{slide_file_title}.pptx"


def generate_ppt(topic, slide_length, api_key):

    client = openai.OpenAI(api_key=api_key)
    msg = message(topic, slide_length)

    response = client.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=[{"role": "user", "content": msg}]
    )

    reply = response.choices[0].message.content
    text_response = parse_response(reply)
    return text_response

def main():

    def button_click():
        # result_label.config(text=generate_ppt('carros populares no Brasil', 5, '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'))
        if input2.get().isdigit():
            result_label.config(text=generate_ppt(input1.get(), input2.get(), input0.get()))

    window = tk.Tk()
    window.title("ChatGPT Generated PPTs!")
    window.configure(padx=100, pady=100)

    label0 = tk.Label(window, text="OpenAI API Key: (Se for inválido, este programa falhará)")
    label0.pack()

    input0 = tk.Entry(window)
    input0.pack()

    label1 = tk.Label(window, text="Escreva-me uma apresentação PPT sobre...")
    label1.pack()

    input1 = tk.Entry(window)
    input1.pack()

    label2 = tk.Label(window, text="Número de slides:")
    label2.pack()

    input2 = tk.Entry(window)
    input2.pack()

    button = tk.Button(window, text="Enviar", command=button_click)
    button.pack()

    result_label = tk.Label(window, text="Resultado")
    result_label.pack()

    window.mainloop()


if __name__ == '__main__':
    main()
